from pptx import Presentation


def extract_pptx_text(file_path: str) -> str:
    presentation = Presentation(file_path)

    lines = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    lines.append(shape.text)

    return "\n".join(lines)